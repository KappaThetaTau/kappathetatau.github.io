import collections 
import collections.abc
from pptx import Presentation
# from pptx.util import Inches
import pandas as pd
import webbrowser

social_fit_form = pd.read_excel(r'C:\Users\lambh\Documents\Github\kappathetatau.github.io\rush-ppt-gen\Spring 2023 Social Fit Form Instructions.xlsx')
# print(social_fit_form.columns.values)
i = 0
prs = Presentation('test_slide.pptx')
prs = Presentation('FA22 Final Voting Slides.pptx')
# print(prs)
# print(prs.shapes())
for slide in prs.slides:
    for shape in slide.shapes:
        try:
            text_box = shape.text
            if (text_box == ''):
                continue
            
            # print("Shape.text:", shape.text)
            
            broken = text_box.split("\n")
            # print(broken)
            if (('Name:' in broken[0])):
                # print("broken:", trim(broken[0][5:len(broken[0])])
                raw_name = broken[0][5:len(broken[0])].strip()
                # print(raw_name)
                print(raw_name)
                continue
            else:
                continue
        except:
            continue
    # 
        # print(shape.text)
        # rush_name = ""
        # if "Name:" in shape.text:
        #     rush_name = shape.text[7:len(shape.text)]
        #     # rush_name = "boo"
        # print(rush_name)
    # except:
        # print("error")
# for link in social_fit_form['Their Slide, screenshot photo for form']:
#     # if (i > 2):
#     #     break
#     # print(link)
#     # webbrowser.open(link)
#     # webbrowser.C
#     i += 1
    # try
    #     webbrowser.open(link)
    # except
# prs = Presentation('FA22 Final Voting Slides.pptx') 
# print("ur mom")
# prs.save('FA22 Final Voting Slides.pptx') 

